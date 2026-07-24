"""PPTX ingest unit testleri (docling → python-pptx fallback)."""

from __future__ import annotations

from unittest.mock import patch

import pytest

from doqqy.ingest.base import IngestError
from doqqy.ingest.pptx_ingest import ingest_pptx
from doqqy.workspace import Workspace


def _fake_pptx(tmp_path):
    path = tmp_path / "deck.pptx"
    path.write_bytes(b"fake pptx content")
    return path


def _real_pptx_with_title_and_body(tmp_path, title: str, body: str):
    """python-pptx ile gerçek bir .pptx dosyası üretir (fallback yolunu test etmek için)."""
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    slide.shapes.title.text = title
    slide.placeholders[1].text_frame.text = body

    path = tmp_path / "deck.pptx"
    prs.save(str(path))
    return path


def test_pptx_ingest_docling_success(tmp_path):
    """Docling başarılı olduğunda doğrudan sonucu dönmeli."""
    ws = Workspace(tmp_path)
    ws.ensure_dirs()
    fake_pptx = _fake_pptx(tmp_path)

    with patch("doqqy.ingest.pptx_ingest._parse_with_docling", return_value="## Slide 1\n\nContent"):
        doc = ingest_pptx(fake_pptx, ws)

        assert doc.content == "## Slide 1\n\nContent"
        assert doc.metadata["parser"] == "docling"
        assert doc.metadata["type"] == "pptx"


def test_pptx_ingest_fallback_slide_titles_become_headings(tmp_path):
    """Docling başarısız olursa python-pptx fallback devreye girmeli; slayt başlığı '##' olmalı."""
    ws = Workspace(tmp_path)
    ws.ensure_dirs()
    real_pptx = _real_pptx_with_title_and_body(tmp_path, "My Slide Title", "First bullet point")

    with patch("doqqy.ingest.pptx_ingest._parse_with_docling", side_effect=RuntimeError("docling boom")):
        doc = ingest_pptx(real_pptx, ws)

        assert doc.metadata["parser"] == "python-pptx"
        assert "## My Slide Title" in doc.content
        assert "First bullet point" in doc.content


def test_pptx_ingest_empty_deck_raises(tmp_path):
    """Boş bir sunum (içerik çıkarılamayan) IngestError fırlatmalı."""
    ws = Workspace(tmp_path)
    ws.ensure_dirs()
    fake_pptx = _fake_pptx(tmp_path)

    with patch("doqqy.ingest.pptx_ingest._parse_with_docling", return_value=""), patch(
        "doqqy.ingest.pptx_ingest._parse_with_python_pptx", return_value=""
    ):
        with pytest.raises(IngestError, match="boş içerik"):
            ingest_pptx(fake_pptx, ws)


def test_pptx_ingest_missing_python_pptx_dep_raises_clear_error(tmp_path):
    """python-pptx kurulu değilse (ImportError) docling hatasıyla birlikte açık bir mesaj vermeli."""
    ws = Workspace(tmp_path)
    ws.ensure_dirs()
    fake_pptx = _fake_pptx(tmp_path)

    with patch("doqqy.ingest.pptx_ingest._parse_with_docling", side_effect=RuntimeError("docling boom")), patch(
        "doqqy.ingest.pptx_ingest._parse_with_python_pptx",
        side_effect=ImportError("No module named 'pptx'"),
    ):
        with pytest.raises(IngestError, match="python-pptx kurulu değil"):
            ingest_pptx(fake_pptx, ws)
