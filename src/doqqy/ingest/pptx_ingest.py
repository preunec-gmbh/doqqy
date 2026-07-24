"""PPTX ingester'ı: docling (ana) → python-pptx (fallback)."""

from __future__ import annotations

from functools import lru_cache
from pathlib import Path
from typing import Any

from doqqy.config import get_logger
from doqqy.ingest.base import Document, IngestError, base_metadata, content_hash, processed_path_for
from doqqy.workspace import Workspace

_LOG = get_logger("doqqy.ingest.pptx")


@lru_cache(maxsize=1)
def _get_docling_converter():
    """Standart Docling converter'ı tek seferlik belleğe yükler."""
    from docling.document_converter import DocumentConverter  # type: ignore

    return DocumentConverter()


def _parse_with_docling(source: Path) -> str:
    converter = _get_docling_converter()
    result = converter.convert(str(source))
    return result.document.export_to_markdown()


def _parse_with_python_pptx(source: Path) -> str:
    """python-pptx → markdown. Slayt başlığı '##', gövde metni paragraf olur."""
    from pptx import Presentation  # type: ignore

    prs = Presentation(str(source))
    sections: list[str] = []

    for idx, slide in enumerate(prs.slides, start=1):
        title_shape = slide.shapes.title
        title_id = title_shape.shape_id if title_shape is not None else None
        title_text = title_shape.text.strip() if title_shape is not None and title_shape.has_text_frame else ""
        heading = title_text or f"Slide {idx}"

        body_paragraphs: list[str] = []
        for shape in slide.shapes:
            if shape.shape_id == title_id or not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text.strip()
                if text:
                    body_paragraphs.append(text)

        section = f"## {heading}"
        if body_paragraphs:
            section += "\n\n" + "\n\n".join(body_paragraphs)
        sections.append(section)

    return "\n\n".join(sections)


def ingest_pptx(source: Path, ws: Workspace, **_kwargs: Any) -> Document:
    md: str | None = None
    parser_used: str | None = None
    docling_error: Exception | None = None

    try:
        md = _parse_with_docling(source)
        parser_used = "docling"
    except Exception as exc:  # noqa: BLE001
        docling_error = exc
        _LOG.warning("docling başarısız (%s): %s", source.name, exc)

    if md is None or not md.strip():
        try:
            fallback_md = _parse_with_python_pptx(source)
            if fallback_md.strip():
                md = fallback_md
                parser_used = "python-pptx"
        except ImportError as exc:
            raise IngestError(
                f"docling başarısız ({docling_error}) ve python-pptx kurulu değil: "
                "'pip install python-pptx' ile kurun."
            ) from exc
        except Exception as exc:  # noqa: BLE001
            raise IngestError(f"hem docling hem python-pptx başarısız. docling: {docling_error}. python-pptx: {exc}") from exc

    if md is None or not md.strip():
        raise IngestError("boş içerik")

    meta = base_metadata(source, ws.root, kind="pptx")
    meta["parser"] = parser_used
    meta["content_hash"] = content_hash(md)

    return Document(
        source_path=source,
        processed_path=processed_path_for(source, ws),
        content=md,
        metadata=meta,
    )
